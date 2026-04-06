"""Populate team slide table with metrics from YAML."""

from pathlib import Path
from datetime import datetime
import json

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def load_team_metrics_yaml(yaml_path):
    """Load the YAML metrics file."""
    yaml_path = Path(yaml_path)
    if not yaml_path.exists():
        raise FileNotFoundError(f"YAML file not found: {yaml_path}")

    with yaml_path.open('r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def load_color_rules(config_path):
    """Load the metric color rules configuration."""
    config_path = Path(config_path)
    if not config_path.exists():
        print(f"Warning: Color rules config not found at {config_path}, using default rules")
        return {}
    
    with config_path.open('r', encoding='utf-8') as f:
        config = json.load(f)
        return config.get('metric_color_rules', {})


def extract_action_plans(team_data, team_name):
    """Extract Action Plans for a specific team."""
    if team_name not in team_data:
        return {}
    
    team_months = team_data[team_name]
    action_plans = {}
    
    # Check each month for Action Plans, starting from most recent
    for month in reversed(list(team_months.keys())):
        metrics = team_months[month]
        if isinstance(metrics, list):
            for metric in metrics:
                metric_name = metric.get('metric', '')
                if metric_name.startswith('Action Plan'):
                    plan_num = metric_name.split()[-1]  # Get the number
                    value = metric.get('value')
                    # Only include non-null, non-empty values
                    if value is not None and str(value).strip() and plan_num not in action_plans:
                        action_plans[plan_num] = str(value).strip()
    
    return action_plans


def get_team_metrics(team_data, team_name):
    """Extract metrics for a specific team."""
    for team_key, team_value in team_data.items():
        if team_key.strip().lower() == team_name.strip().lower():
            return team_value
    raise KeyError(f"Team '{team_name}' not found in YAML data.")


def get_latest_months(team_metrics, count=2):
    """Return the latest month labels from ordered team metrics with actual values."""
    if not team_metrics:
        return []

    month_labels = list(team_metrics.keys())
    selected = []
    for month in reversed(month_labels):
        month_data = team_metrics.get(month, [])
        if any(entry.get('value') is not None for entry in month_data):
            selected.append(month)
            if len(selected) == count:
                break

    if len(selected) == count:
        return list(reversed(selected))

    # Fallback to last available month keys when not enough filled months exist.
    if len(month_labels) < count:
        return month_labels
    return month_labels[-count:]


def find_team_slide(prs, team_name):
    """Find the slide index for a given team by matching text shapes."""
    team_name_lower = team_name.strip().lower()
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text.strip().lower()
                if team_name_lower in text:
                    return idx
    return None


def extract_metric_values(team_metrics, metric_name, months=None):
    """Extract values for a specific metric across months."""
    if months is None:
        months = ['February', 'March']

    values = {}
    for month in months:
        if month in team_metrics:
            for entry in team_metrics[month]:
                if entry.get('metric', '').strip().lower() == metric_name.strip().lower():
                    values[month] = entry.get('value')
                    break
    return values


def format_metric_value(value, metric_name):
    """Format a metric value for display."""
    if value is None:
        return 'N/A'

    if isinstance(value, bool):
        return str(value)

    if isinstance(value, (int, float)):
        # Check if this should be a percentage
        if any(keyword in metric_name.lower() for keyword in ['%', 'rate', 'predictability', 'commitment', 'allocation']):
            return f"{round(value * 100, 1)}%"
        # Check if this is a time-based metric
        if 'days' in metric_name.lower():
            return f"{value} days"
        if 'hours' in metric_name.lower():
            return f"{value} hrs"
        return str(value)

    return str(value)


def calculate_mom_delta(feb_val, mar_val, metric_name=''):
    """Calculate month-over-month delta with appropriate units."""
    if feb_val is None or mar_val is None:
        return 'N/A'

    if isinstance(feb_val, (int, float)) and isinstance(mar_val, (int, float)):
        # Check if this is a percentage metric
        is_percentage = any(keyword in metric_name.lower() for keyword in ['%', 'rate', 'predictability', 'commitment', 'allocation'])
        multiplier = 100 if is_percentage else 1

        delta = (mar_val - feb_val) * multiplier
        delta_rounded = round(delta, 2)

        # Format with arrows instead of +/-
        if delta_rounded > 0:
            formatted_delta = f"↑{delta_rounded}"
        elif delta_rounded < 0:
            formatted_delta = f"↓{abs(delta_rounded)}"  # Use abs to remove the negative sign
        else:
            formatted_delta = f"{delta_rounded}"

        # Add units based on metric type
        if is_percentage:
            return f"{formatted_delta}%"
        elif 'days' in metric_name.lower():
            return f"{formatted_delta} days"
        elif 'hours' in metric_name.lower():
            return f"{formatted_delta} hrs"
        else:
            return formatted_delta

    return 'N/A'


def build_metrics_table(team_metrics, months=None):
    """Build a structured metrics table for slide updates."""
    if months is None:
        months = get_latest_months(team_metrics, 2)

    # Fallback to the default template labels if no months are available.
    if len(months) < 2:
        months = ['February', 'March']

    month_a, month_b = months[0], months[1]

    key_metrics = [
        'Sprint Predictability',
        'Cycle Time (Days)',
        'Unplanned work %',
        'Wasted Effort %',
        'Severity 1 Incidents',
        'Severity 2 Incidents', 
        'Severity 3 Incidents',
        'MTTR (Sev 1)',
        'Dispositioned PDR Bugs - Created',
        'Dispositioned PDR Bugs - Resolved',
        'Change Failure Rate',
        'Code Coverage %',
        'Commitment %',
        'Investment Allocations - R',
        'Investment Allocations - T',
        'Investment Allocations - M',
        'Investment Allocations - 0',
        'Deployment Frequency',
        'Lead Time for Changes',
    ]

    table_data = []

    for metric in key_metrics:
        values = extract_metric_values(team_metrics, metric, [month_a, month_b])
        month_a_val = values.get(month_a)
        month_b_val = values.get(month_b)

        row = {
            'metric': metric,
            'period_a': format_metric_value(month_a_val, metric),
            'period_b': format_metric_value(month_b_val, metric),
            'delta': calculate_mom_delta(month_a_val, month_b_val, metric),
        }
        table_data.append(row)

    # Add combined metrics for rows that aggregate multiple values
    # Sev 1/2/3 Incidents combined row
    sev1_a = extract_metric_values(team_metrics, 'Severity 1 Incidents', [month_a]).get(month_a)
    sev2_a = extract_metric_values(team_metrics, 'Severity 2 Incidents', [month_a]).get(month_a)
    sev3_a = extract_metric_values(team_metrics, 'Severity 3 Incidents', [month_a]).get(month_a)
    
    sev1_b = extract_metric_values(team_metrics, 'Severity 1 Incidents', [month_b]).get(month_b)
    sev2_b = extract_metric_values(team_metrics, 'Severity 2 Incidents', [month_b]).get(month_b)
    sev3_b = extract_metric_values(team_metrics, 'Severity 3 Incidents', [month_b]).get(month_b)
    
    combined_sev_row = {
        'metric': 'Sev 1/2/3 Incidents',
        'period_a': f"{format_metric_value(sev1_a, 'Severity 1 Incidents')}/{format_metric_value(sev2_a, 'Severity 2 Incidents')}/{format_metric_value(sev3_a, 'Severity 3 Incidents')}",
        'period_b': f"{format_metric_value(sev1_b, 'Severity 1 Incidents')}/{format_metric_value(sev2_b, 'Severity 2 Incidents')}/{format_metric_value(sev3_b, 'Severity 3 Incidents')}",
        'delta': 'N/A',  # Combined delta not meaningful
    }
    table_data.append(combined_sev_row)

    # Dispositioned PDR Bugs combined row
    bugs_created_a = extract_metric_values(team_metrics, 'Dispositioned PDR Bugs - Created', [month_a]).get(month_a)
    bugs_resolved_a = extract_metric_values(team_metrics, 'Dispositioned PDR Bugs - Resolved', [month_a]).get(month_a)
    
    bugs_created_b = extract_metric_values(team_metrics, 'Dispositioned PDR Bugs - Created', [month_b]).get(month_b)
    bugs_resolved_b = extract_metric_values(team_metrics, 'Dispositioned PDR Bugs - Resolved', [month_b]).get(month_b)
    
    combined_bugs_row = {
        'metric': 'Dispositioned PDR Bugs',
        'period_a': f"Assigned: {format_metric_value(bugs_created_a, 'Dispositioned PDR Bugs - Created')}\nResolved: {format_metric_value(bugs_resolved_a, 'Dispositioned PDR Bugs - Resolved')}",
        'period_b': f"Assigned: {format_metric_value(bugs_created_b, 'Dispositioned PDR Bugs - Created')}\nResolved: {format_metric_value(bugs_resolved_b, 'Dispositioned PDR Bugs - Resolved')}",
        'delta': f"Assigned: {calculate_mom_delta(bugs_created_a, bugs_created_b, 'Dispositioned PDR Bugs - Created')}\nResolved: {calculate_mom_delta(bugs_resolved_a, bugs_resolved_b, 'Dispositioned PDR Bugs - Resolved')}",
    }
    table_data.append(combined_bugs_row)

    return table_data


def export_metrics_table_csv(table_data, output_path, months=None):
    """Export metrics table to CSV for easy review."""
    if months is None or len(months) < 2:
        months = ['February', 'March']

    output_path = Path(output_path)
    with output_path.open('w', encoding='utf-8', newline='') as f:
        f.write(f'Metric,{months[0]},{months[1]},MoM Delta\n')
        for row in table_data:
            f.write(f'"{row["metric"]}",{row["period_a"]},{row["period_b"]},{row["delta"]}\n')

    return output_path


def update_slide_table(prs, slide_idx, table_data, color_rules, months=None):
    """Update a slide's table with new metrics data, preserving section headers."""
    slide = prs.slides[slide_idx - 1]

    table_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_table') and shape.has_table:
            table_shape = shape
            break

    if not table_shape:
        print(f"No table found on slide {slide_idx}")
        return False

    table = table_shape.table

    if months is None or len(months) < 2:
        months = ['February', 'March']

    # Update the header labels to match the selected month names.
    if len(table.rows) > 0 and len(table.rows[0].cells) >= 3:
        header_cells = table.rows[0].cells
        header_cells[1].text = months[0]
        header_cells[2].text = months[1]
        for cell_idx in [1, 2]:
            for paragraph in header_cells[cell_idx].text_frame.paragraphs:
                paragraph.font.size = Pt(9)

    # Define which rows are section headers that should be preserved
    section_header_rows = {1, 6, 9, 13}  # Rows with section titles like "Productivity & Capacity"
    
    # Create mapping from metric names to table rows
    metric_to_row = {
        'Sprint Predictability': 2,
        'Cycle Time (Days)': 3,
        'Unplanned work %': 4,
        'Wasted Effort %': 5,
        'Sev 1/2/3 Incidents': 7,  # Combined severity incidents
        'MTTR (Sev 1)': 8,
        'Dispositioned PDR Bugs': 10,  # Combined created/resolved
        'Change Failure Rate': 11,
        'Code Coverage %': 12,
        'Commitment %': 14,
    }

    # Update only data rows, skip section headers
    for metric_row in table_data:
        metric_name = metric_row['metric']
        if metric_name in metric_to_row:
            row_idx = metric_to_row[metric_name]
            if row_idx < len(table.rows):
                # Skip known section header rows
                if row_idx in section_header_rows:
                    print(f"Skipping section header row {row_idx}: {table.rows[row_idx].cells[0].text}")
                    continue
                
                cells = table.rows[row_idx].cells
                if len(cells) >= 4:
                    # Update the data cells
                    cells[1].text = metric_row['period_a']
                    cells[2].text = metric_row['period_b']
                    cells[3].text = metric_row['delta']
                    
                    # Set font size to 9pt for data cells (period columns and delta)
                    for cell_idx in [1, 2, 3]:  # Data columns
                        if cells[cell_idx].text_frame.paragraphs:
                            for paragraph in cells[cell_idx].text_frame.paragraphs:
                                paragraph.font.size = Pt(9)
                                for run in paragraph.runs:
                                    run.font.size = Pt(9)
                    
                    # Apply color to delta cell based on rules
                    delta_cell = cells[3]
                    if delta_cell.text_frame.paragraphs:
                        delta_text = delta_cell.text.strip()
                        if delta_text and delta_text != 'N/A':
                            # Determine if this is a positive or negative change
                            is_positive = delta_text.startswith('↑')
                            is_negative = delta_text.startswith('↓')
                            
                            # Get color rule for this metric
                            rule = color_rules.get(metric_name, {})
                            positive_is_good = rule.get('positive_is_good', True)  # Default to positive being good
                            
                            # Determine color: green for good changes, red for bad changes
                            if is_positive and positive_is_good:
                                color = RGBColor(0, 128, 0)  # Green
                            elif is_positive and not positive_is_good:
                                color = RGBColor(255, 0, 0)  # Red
                            elif is_negative and positive_is_good:
                                color = RGBColor(255, 0, 0)  # Red (negative when positive is good)
                            elif is_negative and not positive_is_good:
                                color = RGBColor(0, 128, 0)  # Green (negative when positive is bad)
                            else:
                                color = None  # No color for zero changes
                            
                            if color:
                                for paragraph in delta_cell.text_frame.paragraphs:
                                    paragraph.font.color.rgb = color
                                    for run in paragraph.runs:
                                        run.font.color.rgb = color
                    
                    print(f"Updated row {row_idx} ({metric_name}): {months[0]}={metric_row['period_a']}, {months[1]}={metric_row['period_b']}, ∆={metric_row['delta']} (9pt font)")

    return True


def update_action_plan_placeholders(prs, slide_idx, action_plans):
    """Update Action Plan placeholders on the slide."""
    slide = prs.slides[slide_idx - 1]
    
    # Find the shape containing Action Plan placeholders
    placeholder_shape = None
    for shape in slide.shapes:
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            text = shape.text.strip()
            if 'Focus Area Heading: Description' in text:
                placeholder_shape = shape
                break
    
    if not placeholder_shape:
        print('No Action Plan placeholder found on slide')
        return False
    
    # Build the replacement text
    replacement_lines = []
    
    for i in range(1, 4):  # Action Plan 1, 2, 3
        plan_key = str(i)
        if plan_key in action_plans and action_plans[plan_key] is not None:
            # Replace placeholder with actual content
            plan_content = str(action_plans[plan_key]).strip()
            replacement_lines.append(plan_content)
        else:
            # Remove placeholder (leave empty)
            replacement_lines.append('')
    
    # Join the lines, filtering out empty ones
    final_text = '\n\n'.join(line for line in replacement_lines if line.strip())
    
    # Update the shape text
    placeholder_shape.text = final_text

    # Enforce 9pt font and bold the heading for each action plan paragraph
    if placeholder_shape.text_frame.paragraphs:
        for paragraph in placeholder_shape.text_frame.paragraphs:
            text = paragraph.text
            if ':' in text:
                heading_end = text.index(':') + 1
                heading_text = text[:heading_end]
                remainder_text = text[heading_end:]

                paragraph.text = heading_text
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.bold = True

                if remainder_text:
                    remainder_run = paragraph.add_run()
                    remainder_run.text = remainder_text
                    remainder_run.font.size = Pt(9)
                    remainder_run.font.bold = False
            else:
                paragraph.font.size = Pt(9)
                for run in paragraph.runs:
                    run.font.size = Pt(9)
                    run.font.bold = False

    print(f'Updated Action Plan placeholders: {len([line for line in replacement_lines if line.strip()])} plans added, {len([line for line in replacement_lines if not line.strip()])} placeholders removed')
    
    return True


def main():
    """Main execution."""
    import argparse

    parser = argparse.ArgumentParser(description='Populate a team slide from YAML metrics.')
    parser.add_argument('--team', default='Team Name', help='Team name to populate')
    parser.add_argument('--yaml', default='output/team_metrics.yaml', help='Path to the team metrics YAML file')
    parser.add_argument('--pptx', default='input/March 2026 P&T Scorecards_Talent.pptx', help='Path to the PowerPoint file')
    parser.add_argument('--output', default='output/Team_Updated.pptx', help='Output PowerPoint file')
    parser.add_argument('--csv', default='output/team_metrics.csv', help='CSV export path')
    parser.add_argument('--rules', default='config/metric_color_rules.json', help='Color rules JSON path')
    args = parser.parse_args()

    yaml_path = Path(args.yaml)
    pptx_path = Path(args.pptx)
    output_pptx_path = Path(args.output)
    csv_path = Path(args.csv)
    config_path = Path(args.rules)
    team_name = args.team

    print(f"Loading team metrics from YAML for '{team_name}'...")
    team_data = load_team_metrics_yaml(yaml_path)
    team_metrics = get_team_metrics(team_data, team_name)

    print('Loading color rules configuration...')
    color_rules = load_color_rules(config_path)
    print(f'Loaded color rules for {len(color_rules)} metrics')

    print('Extracting Action Plans...')
    action_plans = extract_action_plans(team_data, team_name)
    print(f'Found {len(action_plans)} Action Plans with content')

    print('Building metrics table...')
    months = get_latest_months(team_metrics, 2)
    table_data = build_metrics_table(team_metrics, months)

    print('\nExtracted metrics (sample):')
    for row in table_data[:5]:
        print(f"  {row['metric']}: {months[0]}={row['period_a']}, {months[1]}={row['period_b']}, ∆={row['delta']}")

    print(f'\nExporting metrics to CSV: {csv_path}')
    export_metrics_table_csv(table_data, csv_path, months)

    print(f'\nLoading PowerPoint: {pptx_path}')
    prs = Presentation(pptx_path)

    slide_idx = find_team_slide(prs, team_name)
    if slide_idx is None:
        print(f"Warning: Could not find a slide for '{team_name}'. Defaulting to slide 3.")
        slide_idx = 3

    print(f"Updating Slide {slide_idx} table with {team_name} metrics...")
    success = update_slide_table(prs, slide_idx=slide_idx, table_data=table_data, color_rules=color_rules, months=months)

    if success:
        print('Updating Action Plan placeholders...')
        update_action_plan_placeholders(prs, slide_idx, action_plans)

    if success:
        print(f'Saving updated presentation to: {output_pptx_path}')
        prs.save(output_pptx_path)
        print(f'✓ Successfully saved to {output_pptx_path}')
    else:
        print('⚠ Slide update failed; presentation not saved.')


if __name__ == '__main__':
    main()
