"""Map team metrics from YAML to the PowerPoint deck."""

from pathlib import Path
import re

import yaml
from pptx import Presentation


def load_team_metrics_yaml(yaml_path):
    yaml_path = Path(yaml_path)
    if not yaml_path.exists():
        raise FileNotFoundError(f"YAML file not found: {yaml_path}")

    with yaml_path.open('r', encoding='utf-8') as f:
        return yaml.safe_load(f)


def find_team_metrics(team_data, team_name):
    for team_key, team_value in team_data.items():
        if team_key.strip().lower() == team_name.strip().lower():
            return team_value
    raise KeyError(f"Team '{team_name}' not found in YAML data.")


def extract_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text.strip()
            if text:
                texts.append(text)

        if hasattr(shape, 'has_table') and shape.has_table:
            table_texts = []
            for row in shape.table.rows:
                row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_data:
                    table_texts.append(' | '.join(row_data))
            if table_texts:
                texts.append('TABLE: ' + ' ; '.join(table_texts))

    return '\n'.join(texts)


def inspect_presentation(pptx_path):
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")

    prs = Presentation(pptx_path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = None
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()

        text = extract_slide_text(slide)
        slides.append({'index': idx, 'title': title, 'text': text})

    return slides


def find_compensation_slide(slides, team_name):
    matches = []
    pattern = re.compile(re.escape(team_name), re.IGNORECASE)
    for slide in slides:
        title = slide['title'] or ''
        text = slide['text'] or ''
        if pattern.search(title) or pattern.search(text):
            matches.append(slide)
    return matches


def print_slide_summary(slide):
    print(f"Slide {slide['index']}")
    print(f"Title: {slide['title'] or '<none>'}")
    print("Text:")
    print(slide['text'])
    print('-' * 80)


def normalize_metric_label(label):
    if not label:
        return label
    return re.sub(r'\s+', ' ', label).strip()


def summarize_compensation_metrics(team_metrics):
    summary = {}
    for month, entries in team_metrics.items():
        summary[month] = [normalize_metric_label(entry.get('metric')) for entry in entries]
    return summary


def main():
    yaml_path = Path('output/team_metrics.yaml')
    pptx_path = Path('input/March 2026 P&T Scorecards_Talent.pptx')

    print('Loading team metrics from YAML...')
    team_data = load_team_metrics_yaml(yaml_path)
    comp_metrics = find_team_metrics(team_data)
    comp_summary = summarize_compensation_metrics(comp_metrics)

    print(f"Found team metrics across {len(comp_summary)} months.")
    for month, metrics in comp_summary.items():
        print(f"\n{month} ({len(metrics)} metrics):")
        for metric in metrics:
            print(f"  - {metric}")

    print('\nInspecting PowerPoint slides...')
    slides = inspect_presentation(pptx_path)
    matches = find_compensation_slide(slides)

    if not matches:
        print('No slides directly mention the team.')
        print('Showing all slides for manual review:')
        for slide in slides:
            print_slide_summary(slide)
        return

    print(f'Found {len(matches)} matching slide(s):')
    for slide in matches:
        print_slide_summary(slide)


if __name__ == '__main__':
    main()
